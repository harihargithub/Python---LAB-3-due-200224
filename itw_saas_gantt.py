from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt

# Create a new presentation
prs = Presentation()

# Slide 1: Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "IT Team Action Plan for Client Portal Deployment"
subtitle.text = "Structured Tasks with Dependencies & Client Inputs"

# Slide 2: Objective
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Objective"
slide.placeholders[1].text = (
    "To launch a fully branded, functional client portal with real-time "
    "flight/hotel bookings, integrated payments, and admin management."
)

# Slide 3: Task Breakdown Table
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Task Breakdown (Part 1)"

# Add a table to the slide
rows = 6
cols = 6
left = Inches(0.3)
top = Inches(1.5)
width = Inches(9)
height = Inches(0.8)
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Set column headers
headers = ["S.No", "Task", "Details", "Client Input", "Mandatory", "Notes"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.text_frame.paragraphs[0].font.bold = True

# Add rows
data = [
    [
        "1",
        "Setup Company Branding",
        "Apply company name across platform",
        "Yes",
        "Yes",
        "Client-provided name",
    ],
    [
        "2",
        "Configure Custom Domain",
        "Setup DNS to point portal",
        "Yes",
        "Yes",
        "Client must update A record",
    ],
    [
        "2a",
        "DNS Setup Help",
        "Provide tutorial for DNS update",
        "No",
        "No",
        "Video link can be shared",
    ],
    [
        "3",
        "Verify DNS Configuration",
        "Confirm A record maps IP",
        "Yes",
        "Yes",
        "Validate before proceeding",
    ],
    [
        "4",
        "Apply Logo",
        "Use client image in templates",
        "Yes",
        "Yes",
        "PNG/JPEG format",
    ],
]

for row_idx, row_data in enumerate(data, start=1):
    for col_idx, val in enumerate(row_data):
        table.cell(row_idx, col_idx).text = val

# Slide 4: Responsibility Assignment
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Responsibility Assignment"
slide.placeholders[1].text = (
    "- Developer: API Integration, Custom Logic\n"
    "- Designer: Logo, Theme, Layout\n"
    "- DevOps: DNS, Domain Setup\n"
    "- Client: Branding Inputs, Content Pages"
)

# Slide 5: Client Collaboration Points
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Client Collaboration Points"
slide.placeholders[1].text = (
    "Tasks that require client input for smooth deployment:\n\n"
    "- Company Branding\n"
    "- Domain Setup and DNS Update\n"
    "- Logo and Color Theme\n"
    "- Content for Static Pages\n"
    "- API Credentials for Flights, Hotels, Payments"
)

# Slide 6: Summary and Next Steps
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Summary & Next Steps"
slide.placeholders[1].text = (
    "✔️ Internal Testing\n"
    "✔️ Client UAT (User Acceptance Testing)\n"
    "✔️ Go-Live with branding & real-time APIs\n"
    "✔️ Optional Phase 2: Add Loyalty Points, Chatbot, CRM"
)

# Slide 7: Gantt Chart
# Create Gantt chart using matplotlib
tasks = [
    ("Setup Company Branding", 0, 2),
    ("Configure Custom Domain", 2, 4),
    ("Verify Domain DNS Configuration", 4, 5),
    ("Upload and Apply Logo", 5, 6),
    ("Upload and Apply Favicon", 6, 7),
    ("Apply Primary & Secondary Colors", 7, 8),
    ("Apply Template Layout", 8, 9),
    ("Integrate Flight APIs", 9, 11),
    ("Integrate Client Flight API", 11, 12),
    ("Integrate Client Hotel API", 12, 13),
    ("Setup TTW Payment Gateway", 13, 14),
    ("Integrate Client Payment Gateway", 14, 15),
    ("Configure Business Email Notifications", 15, 16),
    ("Configure Footer Info", 16, 17),
    ("Create Static Pages", 17, 18),
    ("Set Language and Currency Defaults", 18, 19),
    ("Configure Database", 19, 20),
    ("Automated email & sms", 20, 21),
    ("Client Admin Panel", 21, 22),
]

fig, ax = plt.subplots(figsize=(10, 6))
for i, (task, start, end) in enumerate(tasks):
    ax.barh(i, end - start, left=start, height=0.6, align="center", color="skyblue")
    ax.text(start + 0.1, i, task, va="center", ha="left", fontsize=8)

ax.set_yticks(range(len(tasks)))
ax.set_yticklabels([""] * len(tasks))  # Hide task names on y-axis, shown on bars
ax.set_xlabel("Days")
ax.set_title("Gantt Chart for IT Team Portal Deployment")
# Remove top and right spines(lines)
ax.spines["top"].set_visible(False)
ax.spines["right"].set_visible(False)
plt.tight_layout()
gantt_path = r"C:\Users\nhari\Documents\gantt_chart.png"  # Save to a valid directory
plt.savefig(gantt_path)
plt.close()

# Add Gantt chart slide
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Visual Timeline – Gantt Chart"

# Insert Gantt chart image
left = Inches(1)
top = Inches(1.5)
height = Inches(4.5)
slide.shapes.add_picture(gantt_path, left, top, height=height)

# Save the presentation
pptx_path = r"C:\Users\nhari\OneDrive\Documents\itw\b2bta\saas\IT_Team_Client_Portal_Action_Plan_Gantt.pptx"
prs.save(pptx_path)

print(f"Presentation saved to {pptx_path}")
