from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create a new presentation
prs = Presentation()

# Define title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "IT Team Action Plan for Client Portal Deployment"
subtitle.text = "Structured Tasks with Dependencies & Client Inputs"

# Slide 2: Objective
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Objective"
slide.placeholders[1].text = (
    "To launch a fully branded, functional client portal with real-time "
    "flight/hotel bookings, integrated payments, and admin management."
)

# Slide 3: Task Breakdown Table
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Task Breakdown (Part 1)"

# Add a table to the slide
rows = 6
cols = 6
left = Inches(0.3)
top = Inches(1.5)
width = Inches(9)
height = Inches(0.8)
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Set column headers
headers = ["S.No", "Task", "Details", "Client Input", "Mandatory", "Notes"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.text_frame.paragraphs[0].font.bold = True

# Add rows
data = [
    [
        "1",
        "Setup Company Branding",
        "Apply company name across platform",
        "Yes",
        "Yes",
        "Client-provided name",
    ],
    [
        "2",
        "Configure Custom Domain",
        "Setup DNS to point portal",
        "Yes",
        "Yes",
        "Client must update A record",
    ],
    [
        "2a",
        "DNS Setup Help",
        "Provide tutorial for DNS update",
        "No",
        "No",
        "Video link can be shared",
    ],
    [
        "3",
        "Verify DNS Configuration",
        "Confirm A record maps IP",
        "Yes",
        "Yes",
        "Validate before proceeding",
    ],
    [
        "4",
        "Apply Logo",
        "Use client image in templates",
        "Yes",
        "Yes",
        "PNG/JPEG format",
    ],
]

for row_idx, row_data in enumerate(data, start=1):
    for col_idx, val in enumerate(row_data):
        table.cell(row_idx, col_idx).text = val

# Slide 4: Responsibility Assignment
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Responsibility Assignment"
slide.placeholders[1].text = (
    "- Developer: API Integration, Custom Logic\n"
    "- Designer: Logo, Theme, Layout\n"
    "- DevOps: DNS, Domain Setup\n"
    "- Client: Branding Inputs, Content Pages"
)

# Slide 5: Gantt Chart
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Gantt Chart for IT Team Portal Deployment"

# Define Gantt chart data
tasks = [
    {"name": "Setup Company Branding", "start": 0, "end": 2},
    {"name": "Configure Custom-Domain", "start": 2, "end": 4},
    {"name": "Verify Domain DNS Configuration", "start": 4, "end": 6},
    {"name": "Upload and Apply Logo", "start": 6, "end": 8},
    {"name": "Upload and Apply Favicon", "start": 8, "end": 10},
    {"name": "Apply Primary & Secondary Colors", "start": 10, "end": 12},
    {"name": "Apply Template Layout", "start": 12, "end": 14},
    {"name": "Integrate Flight-APIs", "start": 14, "end": 16},
    {"name": "Integrate Client Flight API", "start": 16, "end": 18},
    {"name": "Integrate Client Hotel API", "start": 18, "end": 20},
    {"name": "Setup TTW Payment Gateway", "start": 20, "end": 22},
    {"name": "Integrate Client Payment Gateway", "start": 22, "end": 24},
    {"name": "Configure Business Email-Notifications", "start": 24, "end": 26},
    {"name": "Configure Footer Info", "start": 26, "end": 28},
    {"name": "Create Static Pages", "start": 28, "end": 30},
    {"name": "Set Language and Currency Defaults", "start": 30, "end": 32},
    {"name": "Configure Database", "start": 32, "end": 34},
    {"name": "Automated email & sms", "start": 34, "end": 36},
    {"name": "Client Admin-Panel", "start": 36, "end": 38},
]

# Draw Gantt chart
chart_left = Inches(1)
chart_top = Inches(1.5)

for i, task in enumerate(tasks):
    # Task label
    label_left = chart_left - Inches(1.5)
    label_top = chart_top + Inches(i * 0.4)
    textbox = slide.shapes.add_textbox(label_left, label_top, Inches(2), Inches(0.4))
    textbox.text = task["name"]
    textbox.text_frame.paragraphs[0].font.size = Pt(10)

    # Task bar
    bar_left = chart_left + Inches(task["start"] * 0.2)
    bar_top = chart_top + Inches(i * 0.4)
    bar_width = Inches((task["end"] - task["start"]) * 0.2)
    bar_height = Inches(0.3)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, bar_left, bar_top, bar_width, bar_height
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 102, 204)  # Blue color

# Add timeline labels
timeline_top = chart_top + Inches(len(tasks) * 0.4)
for day in range(0, 41, 5):  # Timeline from 0 to 40 days, step of 5
    day_left = chart_left + Inches(day * 0.2)
    day_textbox = slide.shapes.add_textbox(
        day_left, timeline_top, Inches(0.5), Inches(0.3)
    )
    day_textbox.text = str(day)
    day_textbox.text_frame.paragraphs[0].font.size = Pt(10)

# Slide 6: Summary and Next Steps
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Summary & Next Steps"
slide.placeholders[1].text = (
    "✔️ Internal Testing\n"
    "✔️ Client UAT (User Acceptance Testing)\n"
    "✔️ Go-Live with branding & real-time APIs\n"
    "✔️ Optional Phase 2: Add Loyalty Points, Chatbot, CRM"
)

# Save the presentation
pptx_path = r"C:\Users\nhari\OneDrive\Documents\itw\b2bta\saas\IT_Team_Client_Portal_Action_Plan_Gantt.pptx"
prs.save(pptx_path)

print(f"Presentation saved to {pptx_path}")
